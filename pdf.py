from pptx import Presentation
import json

from js import document, Uint8Array, File, URL
from pyodide.ffi.wrappers import add_event_listener


def downloadPdf(event):
    print("상장 제작중...")
    data = json.loads(document.getElementById("payload").value)
    print(data)

    prs = Presentation("gdsc_deu_23-24.pptx")
    slide = prs.slides[0]

    # name -> data[name] 으로 변경
    for k, v in data.items():
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.text = run.text.replace(k, v)

    file_name = data["name"] + "-" + data["title"] + ".pptx"

    prs.save(file_name)

    # 생성된 문서 다운로드
    with open(file_name, "rb") as file:
        data = file.read()
        blob = Uint8Array.new(data)
        file = File.new([blob], file_name)
        url = URL.createObjectURL(file)
        download_link = document.createElement("a")
        download_link.href = url
        download_link.download = file_name
        download_link.click()


download_btn = document.getElementById("download")

add_event_listener(download_btn, "click", downloadPdf)
